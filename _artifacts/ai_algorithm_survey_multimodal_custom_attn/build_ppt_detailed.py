from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = Path(__file__).resolve().parents[2]
ROOT = REPO / '01_ai_infra/kernel/custom_attn'
OUT = ROOT / 'supplements/multimodal-custom-attention.pptx'
ASSET = ROOT / 'assets/papers'
COSMOS_ASSET = REPO / '02_model_systems/multimodal_generation/assets/papers/cosmos-3'

INK = '17212B'; PAPER = 'F7F5EF'; WHITE = 'FFFFFF'; GREY = '52616B'; MINT = '0D9488'; GOLD = 'D5A33B'; CORAL = 'E4573D'; PALE = 'E7F1EE'; LINE = 'D5DDDA'; LAV = 'EAE6F4'

def c(value): return RGBColor.from_string(value)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title = '多模态稀疏 Attention 与定制 Mask Kernel 调研'
prs.core_properties.subject = '图文精读与内核设计评审'

def add_box(slide, x, y, w, h, fill=WHITE, line=None, rounded=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = c(fill)
    shp.line.color.rgb = c(line or fill)
    return shp

def add_text(slide, content, x, y, w, h, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT, font='Aptos', margin=.04, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin); tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = content; r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = c(color); r.font.bold = bold
    tf.vertical_anchor = valign
    return tb

def add_title(slide, title, dark=False, subtitle='图文精读版 | 2026-07-10 | NVIDIA CUDA 重点'):
    add_text(slide, title, .55, .20, 12.15, .70, 24, WHITE if dark else INK, True, font='Aptos Display')
    add_text(slide, subtitle, .58, .98, 11.8, .18, 8, 'B7C8C4' if dark else GREY)

def add_footer(slide, number, source=None, dark=False):
    if source:
        add_text(slide, source, .58, 6.82, 11.6, .18, 8, 'AAB8B5' if dark else GREY)
    add_text(slide, f'{number:02d}', 12.42, 6.80, .35, .2, 10, 'AAB8B5' if dark else GREY, True, PP_ALIGN.RIGHT)

def fit_image(slide, path, x, y, w, h, line='D6DFDB'):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    dx, dy = x + (w - dw) / 2, y + (h - dh) / 2
    add_box(slide, x, y, w, h, WHITE, line, False)
    slide.shapes.add_picture(str(path), Inches(dx), Inches(dy), width=Inches(dw), height=Inches(dh))

def callout(slide, heading, body, x, y, w, h, accent=MINT):
    add_box(slide, x, y, w, h, WHITE, LINE, True)
    add_box(slide, x, y, .10, h, accent, accent)
    add_text(slide, heading, x+.22, y+.14, w-.34, .26, 13, INK, True)
    add_text(slide, body, x+.22, y+.48, w-.34, h-.58, 11, GREY)

def paper_slide(num, title, image, problem, mechanism, impl, source, accent=MINT, image_box=(.55, 1.35, 7.45, 4.95)):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c(PAPER)
    add_title(slide, title)
    x,y,w,h = image_box; fit_image(slide, image, x,y,w,h)
    callout(slide, '解决的问题', problem, 8.25, 1.38, 4.48, 1.36, accent)
    callout(slide, '图中机制', mechanism, 8.25, 2.95, 4.48, 1.53, GOLD)
    callout(slide, '实现 / kernel 含义', impl, 8.25, 4.70, 4.48, 1.58, CORAL)
    add_footer(slide, num, source)
    return slide

# 1
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(INK)
add_text(s, '多模态稀疏 Attention\n与定制 Mask Kernel', .7, .95, 7.5, 1.45, 34, WHITE, True, font='Aptos Display')
add_text(s, '从论文图示、核心机制到 mask lowering、kernel metadata 与 host-device 数据流', .73, 2.73, 7.6, .4, 16, 'C6D6D2')
for i,(a,b,col) in enumerate([('理解', 'high-resolution VLM token selection', MINT),('统一', 'two-way stream lowering / special causal', GOLD),('生成', 'window, CSR, router, temporal reuse', CORAL)]):
    y=3.90+i*.78; add_box(s,8.72,y,3.82,.58,col,col,True); add_text(s,a,8.94,y+.16,.65,.18,12,INK,True); add_text(s,b,9.62,y+.09,2.72,.38,10,INK,False,valign=MSO_ANCHOR.MIDDLE)
add_text(s, '10 篇代表工作 | 原论文图已嵌入 | 代码证据与推断严格区分', .75, 6.75, 8.2, .25, 11, 'A9BFBA'); add_footer(s,1,dark=True)

# 2
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'如何读本报告：一张图必须落到一个可执行对象')
steps=[('任务与 token 拓扑','text / patch / frame / action / chunk',MINT),('可见性语义','causal / local / anchor / read-only',GOLD),('lowering','varlen call / BlockMask / CSR / pack',CORAL),('kernel 与 runtime','tile traversal / plan / KV / CP',MINT)]
for i,(a,b,col) in enumerate(steps):
    x=.65+i*3.17; add_box(s,x,2.35,2.72,1.55,WHITE,LINE,True); add_box(s,x,2.35,2.72,.12,col,col,False); add_text(s,a,x+.2,2.68,2.25,.25,15,INK,True); add_text(s,b,x+.2,3.2,2.22,.4,11,GREY)
    if i<3: add_text(s,'->',x+2.76,2.94,.35,.25,18,INK,True,PP_ALIGN.CENTER)
add_text(s,'每篇均回答：解决什么问题？图中 token/block 是什么？mask 怎样表示？kernel 到底跳过了什么？哪部分只有论文证据？',.75,4.85,11.8,.5,18,INK,True)
callout(s,'读图规则','论文图说明算法语义；源码路径说明 runtime 真正接收什么；性能图只能支持其报告设置，不能自动归因给单个 kernel。',.85,5.7,11.6,.82,GOLD); add_footer(s,2)

# 3
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'领域图谱：mask 不再只是一个 score bias')
rows=[('规则型','Causal-rCM / Cosmos 3','block schedule, stream split','BlockMask / varlen calls',MINT),('索引型','LVSA','CSR indptr + indices','FlashInfer planning',GOLD),('选择型','VMoBA / TSA / VLM FlexAttn','index set + compact QKV','FlashAttention varlen',CORAL),('替代型','FrameDiT','frame-level matrix representation','不生成 token-pair mask',MINT)]
for i,(a,b,d,e,col) in enumerate(rows):
    y=1.45+i*1.16; add_box(s,.75,y,11.8,.84,WHITE,LINE,True); add_box(s,.75,y,1.48,.84,col,col,True); add_text(s,a,.89,y+.28,1.15,.2,13,INK,True,PP_ALIGN.CENTER); add_text(s,b,2.52,y+.17,2.3,.24,14,INK,True); add_text(s,d,5.06,y+.17,3.05,.4,12,GREY); add_text(s,e,8.6,y+.17,3.52,.4,12,INK,True)
add_text(s,'系统趋势：将“可见性”保留为规则、索引或紧凑序列，直到 kernel plan；禁止在模型脚本侧先 materialize L x L。',.75,6.35,11.6,.3,16,INK,True); add_footer(s,3)

# 4-13 individual papers
paper_slide(4,'理解侧：FlexAttention VLM 只让 query 读取被选中的高分辨率细节', ASSET/'flexattention-vlm/fig2_hierarchical_vlm_selection_caption.png', '全量 high-resolution patch 进入 decoder 会产生 quadratic query/key work；只降采样又丢小文字和小目标。', '低分辨率与文本保留 N-token residual stream；上一层 attention map 选择约 10% high-resolution features，仅作为额外 K/V。', '数学对象是 N x (N+M) rectangular attention。compact/varlen 是实现建议；论文未声称 PyTorch FlexAttention 或特定 sparse CUDA backend。V100 上 TextVQA 总时间比 HD/XAttn 低 13%/24%。', 'FlexAttention Fig.2 & Table 5; ECCV 2024; code f814be5', MINT)
paper_slide(5,'统一模型：Cosmos 3 先拆语义，再复用两类 varlen attention', COSMOS_ASSET/'two-way-attention-infra.png', 'reasoner 必须 causal 且不能读取 noisy generator；generator 又要读取同样本 AR 条件与全部 DM tokens。', '同一 packed sample lower 为 AR causal call 与 DM full rectangular call；跨样本用 offsets 隔离。', '收益属于 semantic lowering + FA3/varlen/runtime 组合，不改变模型候选质量。Nano/Super 训练 MFU 仅 0.23/0.30，说明 attention 之外仍有 loader、VAE、通信与 checkpoint 成本。', 'Cosmos 3 Sec.5, Fig.14/16; cosmos-framework 3a5314b', GOLD)
paper_slide(6,'流式世界模型：Causal-rCM 的 special causal mask', ASSET/'causal-rcm/fig3_causal_training_paradigms_caption.png', 'AR diffusion 的 TF 需要 clean history + noisy target；SF 又需 self-generated rollout 与 KV cache。普通三角 causal mask 不够。', '图中 TF / DF / SF 区分 clean、noisy、self-generated history；noisy block 只读允许的 clean history 与自身 block。', 'BlockPattern + AttnMaskSpec -> Flex BlockMask / range metadata；同一 mask 进入 primal 和 JVP Triton kernel。Magi backward 限制需单独核验。', 'Causal-rCM Fig.3; code commit ed3cb14', CORAL)
paper_slide(7,'长视频：LVSA 用 window + rotating anchors 保持稀疏预算', ASSET/'lvsa/fig1_expanded_window_caption.png', '固定 window 会漏长程依赖；window 与 global anchor 重叠又浪费固定 attention budget。', 'expanded local window 与 periodic global frames 构成 A(t)=G union W(t)，每个 query frame 近似保持相同 attended set 大小。', 'frame-block CSR int32 indptr/indices；FlashInfer BlockSparseAttentionWrapper 跳过未列 tile。CPU planner 留 metadata 在 host，并非 GPU kernel 直接读 CPU RAM。', 'LVSA Fig.1; code commit 1ebcc92', MINT)
paper_slide(8,'学习式 block router：VMoBA 的 partition -> select -> varlen', ASSET/'vmoba/fig2_vmoba_pipeline_caption.png', '一维均匀 MoBA block 与视频 temporal/spatial/3D 邻域不匹配；固定 top-k 也浪费异质 head 预算。', 'recurrent 1D/2D/3D partition，mean key 产生 block score；global/threshold selection 后仅在选中 blocks attention。', 'GPU gate + topk/threshold -> nonzero -> gather QKV -> cu_seqlens -> FlashAttention varlen。控制面为 gate/sort/pack/LSE，不传 CSR 或 pair-mask。', 'VMoBA Fig.2; code commit 48aaccd', GOLD)
paper_slide(9,'动态控制面：HASTE 的 mask reuse 与 head-wise calibration', ASSET/'haste/fig4_tmr_ebc_framework_caption.png', 'Video DiT 要多 step denoise；逐 head、逐 step 重算 sparse mask 可能吃掉 attention 节省，统一 threshold 又不公平。', 'TMR 用 Q/K drift 判断每个 head 复用还是刷新 cached descriptor；EBC 用离线 error curve 分配 head-specific threshold。', '应缓存 sparse descriptor，而非 N x N mask。官方代码未取得，CSR/BlockMask/host-device placement 不可断言；机制与结果为 PDF-only。', 'HASTE Fig.4, arXiv:2605.14513', CORAL)
paper_slide(10,'Sparse VideoGen：spatial / temporal head dispatch', ASSET/'sparse-videogen/fig4_svg_workflow_caption.png', 'Video DiT 3D full attention 成为主成本；直接移植 LLM mask 会丢 temporal dependency。', 'sampled rows 对 spatial / temporal / full attention 做 MSE 近似，按 head dispatch 到专用 pattern。', '关键不只是找 pattern：temporal slash 必须 layout transform 才能有 coalesced tile access。论文称 Triton/FlashInfer prototype，具体 metadata 本次未有源码核验。', 'Sparse VideoGen Fig.4, arXiv:2502.01776', MINT)
paper_slide(11,'Token Sparse Attention：保留 selector 灵活性，复用成熟 kernel', ASSET/'token-sparse-attention/fig3_compress_attention_scatter_caption.png', 'block sparse 粒度高效但 token importance 随 layer/head 变化；永久 drop token 会妨碍后续层重新选择。', 'per-head select token subset；compress Q/K/V；在 compact tensors attention；scatter output back 后叠加 residual。', 'kernel 只见连续 compact QKV，可复用 FlashAttention；真实成本为 selector + gather/contiguous + scatter。实现细节为 PDF-only。', 'Token Sparse Attention Fig.3, arXiv:2602.03216', GOLD)
paper_slide(12,'长上下文桥接：MInference 的 pattern-aware index 与 kernel dispatch', ASSET/'minference/fig3_sparse_patterns_caption.png', 'prefill attention 延迟主导；固定 top-k index 跨 prompt recall 显著下降。', '离线给 head 选 A-shape / vertical-slash / block-sparse family；在线建立具体 ranges/columns/blocks。', 'pattern-specific index 交给 PIT/Triton/FlashAttention 类 kernel。迁移到 video 必须改为双向时空 pattern，并处理每 step 的 planner 成本。', 'MInference Fig.3, NeurIPS 2024 Spotlight', CORAL)
paper_slide(13,'架构替代：FrameDiT 以 matrix attention 改变 temporal topology', ASSET/'framedit/fig1_matrix_attention_architecture_caption.png', 'full 3D attention 表达强但昂贵；local factorized attention 便宜却错过大运动。', '以 frame matrix 为对象做 temporal Matrix Attention，Global-Local hybrid 保留局部路径，不再构造原 token-pair temporal mask。', '公开代码仍把 2D mask 转为 -10000 dense bias 并 broadcast；论文算法收益不自动等于 custom sparse kernel。', 'FrameDiT Fig.1, CVPR 2026 Findings; code commit 359bd12', MINT)

# 14 detailed Causal / LVSA evidence visual
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'两条“真正下沉到 runtime”的路径：规则与 CSR')
fit_image(s, ASSET/'causal-rcm/fig4_recipe_comparison_caption.png', .55,1.35,5.95,3.72)
fit_image(s, ASSET/'lvsa/fig4_wall_time_scaling_caption.png', 6.82,1.35,5.95,3.72)
callout(s,'Causal-rCM：规则型','BlockPattern 让 kernel/BlockMask 按 block id 判定 visibility；同一规则覆盖 TF、JVP、cache。',.7,5.28,5.55,1.02,CORAL)
callout(s,'LVSA：索引型','CSR 与 compact layout 交给 FlashInfer planning；图表说明长 horizon/80GB 情况，但不可跨模型横比。',7.0,5.28,5.55,1.02,MINT)
add_footer(s,14,'Causal-rCM Fig.4; LVSA Fig.4')

# 15 compare representation
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'定制 mask 的四种表达：什么能真正跳过 tile？')
cols=[('Rule / BlockMask','block id, window, stream, offset','Causal-rCM','kernel/compiler 可判定 tile','规则须 block-aligned',MINT),('CSR / page table','indptr, indices, page id','LVSA','scheduler 只遍历 nnz block','plan / locality',GOLD),('Selected segments','indices, cu_seqlens, compact QKV','VMoBA / TSA / VLM','标准 varlen kernel','pack/unpack',CORAL),('Dense bias','bool / -inf score bias','FrameDiT public code','通常不跳 tile','L² / dense work','8C9A9C')]
for i,(h,a,b,d,e,col) in enumerate(cols):
    x=.44+i*3.18; add_box(s,x,1.45,2.88,4.8,WHITE,LINE,True); add_box(s,x,1.45,2.88,.52,col,col,True); add_text(s,h,x+.16,1.59,2.55,.2,12,INK,True); add_text(s,a,x+.18,2.27,2.48,.54,12,GREY); add_text(s,b,x+.18,3.05,2.48,.35,14,INK,True); add_text(s,d,x+.18,3.73,2.48,.48,12,GREY); add_text(s,'成本：'+e,x+.18,5.25,2.42,.3,11,CORAL,True)
add_footer(s,15,'Cross-paper implementation synthesis')

# 16 host device
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'长序列 host-device 数据流：传 metadata，不传 dense pair mask')
blocks=[('CPU static planner','geometry / cache\nCSR int32',MINT),('GPU dynamic selector','Q/K drift / top-k\nindices',GOLD),('Plan / pack','page plan / compact QKV\ncu_seqlens',CORAL),('Attention kernel','only nnz tiles /\ncompact sequence',MINT)]
for i,(a,b,col) in enumerate(blocks):
    x=.6+i*3.18; add_box(s,x,2.15,2.56,1.35,col,col,True); add_text(s,a,x+.16,2.43,2.25,.26,14,INK,True,PP_ALIGN.CENTER); add_text(s,b,x+.16,2.86,2.25,.35,11,INK,False,PP_ALIGN.CENTER)
    if i<3: add_text(s,'->',x+2.62,2.65,.35,.25,18,INK,True,PP_ALIGN.CENTER)
callout(s,'允许','静态 window/anchor：CPU 一次性 CSR、pinned metadata、FlashInfer plan；每 request 的 page list 也可由 host scheduler 供给。',.82,4.25,5.7,1.30,MINT)
callout(s,'禁止','CPU 生成 L x L bool/fp16 mask 再拷 GPU。64K 单 bool mask 已 4GiB；per-step top-k CPU 往返还会同步 pipeline。',6.85,4.25,5.7,1.30,CORAL)
add_footer(s,16,'Direct code evidence: LVSA sparse_attention.py:275-304, 598-607')

# 17 blueprint
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(PAPER); add_title(s,'建议的实现边界：先 lowering，再选 kernel')
add_box(s,.7,1.35,12.0,.65,INK,INK,True); add_text(s,'MaskSemantics: stream partition + temporal geometry + dynamic source',.95,1.57,11.5,.2,15,WHITE,True,PP_ALIGN.CENTER)
for i,(h,b,col) in enumerate([('A. rectangles','split into causal/full varlen calls',MINT),('B. regular graph','BlockMask / predicate',GOLD),('C. explicit graph','CSR/page metadata + plan',CORAL),('D. dynamic selection','gather + compact varlen',MINT)]):
    x=.7+i*3.02; add_box(s,x,2.52,2.7,1.25,WHITE,LINE,True); add_text(s,h,x+.15,2.78,2.35,.25,14,INK,True); add_text(s,b,x+.15,3.19,2.35,.28,11,GREY)
add_text(s,'build_attention_plan(spec, qkv_layout, device) -> Plan\nattention_run(q, k, v, Plan) -> out',.98,4.43,11.0,.65,18,INK,True,PP_ALIGN.CENTER,font='Consolas')
callout(s,'评测公式','Ttotal = Tselect + Tmetadata + TH2D/plan + Tpack + Tattn + Tunpack。只报 attention FLOPs 会遗漏控制面和数据搬运。',.95,5.43,11.4,.90,GOLD); add_footer(s,17)

# 18 checklist / conclusion
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=c(INK); add_title(s,'结论：让 mask 的语义在正确层级消失或变小',True)
items=[('1','能拆就拆','双流 / 矩形可见性 -> causal/full varlen calls',MINT),('2','需稀疏就索引化','window/anchor -> CSR / BlockMask / page table',GOLD),('3','需动态就紧凑化','router/selector -> GPU indices + compact QKV',CORAL),('4','控制面也是性能','planner、H2D、top-k、pack/unpack 必须和 attention 同测',MINT)]
for i,(n,h,b,col) in enumerate(items):
    y=1.45+i*1.08; add_box(s,.75,y,.52,.52,col,col,True); add_text(s,n,.93,y+.12,.18,.16,13,INK,True,PP_ALIGN.CENTER); add_text(s,h,1.65,y+.06,2.4,.24,17,WHITE,True); add_text(s,b,4.2,y+.11,7.5,.25,14,'C8D6D3')
add_text(s,'完整图文精读、代码 commit、图清单与 QA：Custom Attention README -> Survey -> Paper -> Asset',.78,6.7,11.8,.22,10,'A9BFBA'); add_footer(s,18,dark=True)

prs.save(OUT)
print(OUT)
