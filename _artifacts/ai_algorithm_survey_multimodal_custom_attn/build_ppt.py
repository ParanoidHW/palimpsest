from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = '01_ai_infra/kernel/custom_attn/多模态稀疏Attention与定制Mask_Kernel调研.pptx'
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
prs.core_properties.title = '多模态稀疏 Attention 与定制 Mask Kernel 调研'

INK = '17212B'; MINT = '0D9488'; CORAL = 'E4573D'; GOLD = 'D5A33B'; PAPER = 'F7F5EF'; GREY = '52616B'; PALE = 'E7F1EE'; WHITE = 'FFFFFF'

def rgb(s): return RGBColor.from_string(s)
def box(slide, x,y,w,h, fill=WHITE, line=None, radius=False):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb=rgb(fill)
    shape.line.color.rgb=rgb(line or fill)
    return shape
def text(slide, s,x,y,w,h, size=16,color=INK,bold=False,align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(.06); tf.margin_top=Inches(.03)
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=s; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(color); tf.vertical_anchor=valign
    return tb
def title(slide, s, dark=False):
    text(slide,s,.55,.32,12.1,.5,28,WHITE if dark else INK,True,font='Aptos Display')
    text(slide,'2026-07-10 | NVIDIA CUDA | 证据包：_artifacts/ai_algorithm_survey_multimodal_custom_attn',.58,.9,12,.24,9,'B9C6C8' if dark else GREY)
def foot(slide,n,dark=False): text(slide,f'{n:02d}',12.55,7.05,.35,.2,10,'A5B2B2' if dark else GREY,True,align=PP_ALIGN.RIGHT)
def card(slide, hdr, body,x,y,w,h,accent=MINT):
    box(slide,x,y,w,h,WHITE,'D7DEDC',True); box(slide,x,y,.09,h,accent,accent)
    text(slide,hdr,x+.22,y+.16,w-.38,.28,15,INK,True)
    text(slide,body,x+.22,y+.55,w-.38,h-.65,12,GREY)

# 1 title
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(INK)
text(s,'多模态稀疏 Attention\n与定制 Mask Kernel',.65,1.05,8.2,1.6,35,WHITE,True,font='Aptos Display')
text(s,'从 mask 语义、稀疏表征到 kernel / planner / KV runtime 的实现趋势',.7,2.9,7.8,.45,16,'C9D8D5')
for i,(a,b,c) in enumerate([('规则','BlockMask / predicate',MINT),('索引','CSR / page table',GOLD),('打包','selector -> varlen',CORAL)]):
    box(s,8.9,1.15+i*1.28,3.55,.92,c,c,True); text(s,a,9.15,1.33+i*1.28,.9,.24,14,INK,True); text(s,b,10.12,1.33+i*1.28,2.1,.3,13,INK)
text(s,'内核设计评审版 | 9 篇深读 | CUDA 重点',.7,6.62,6,.3,12,'A9BFBA'); foot(s,1,True)

# 2 exec
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'核心判断：不要把可见性先展开成 L x L 张量')
text(s,'完整 dense mask 既浪费内存，也不保证跳过 QK/softmax/AV tile。正确的系统边界是：语义 lowering -> 稀疏 metadata / compact QKV -> kernel plan -> attention run。',.62,1.28,12.1,.48,17,INK,True)
for i,(n,h,b,c) in enumerate([('01','能拆矩形','多次 causal/full varlen 调用',MINT),('02','结构化图','CSR / BlockMask / page table',GOLD),('03','动态选择','indices + pack + varlen attention',CORAL)]):
    x=.75+i*4.15; box(s,x,2.25,3.65,2.55,WHITE,'D4D9D6',True); text(s,n,x+.25,2.5,.6,.3,24,c,True); text(s,h,x+.25,3.05,2.8,.28,18,INK,True); text(s,b,x+.25,3.55,2.95,.65,14,GREY)
text(s,'结论：稀疏的单位应该是 kernel 能跳过的 tile、page 或 compact segment，而非抽象的 0/1 score bias。',.78,5.55,11.7,.45,18,INK,True); foot(s,2)

# 3 timeline
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'时间线：从 kernel-aware pattern 到多模态 runtime')
box(s,.85,3.55,11.65,.07,INK,INK)
items=[('2024','MInference','per-head dynamic pattern',MINT),('2025','Sparse VideoGen','spatial / temporal heads',GOLD),('2025','VMoBA','block router + varlen',CORAL),('2026','HASTE / LVSA','reuse / CSR + FlashInfer',MINT),('2026','Causal-rCM / Cosmos 3','custom JVP / two-way lowering',CORAL)]
for i,(yr,n,b,c) in enumerate(items):
    x=.82+i*2.42; box(s,x,2.57,.18,.18,c,c,True); text(s,yr,x,1.62,2.05,.3,17,INK,True); text(s,n,x,2.03,2.05,.38,14,INK,True); text(s,b,x,3.92,2.02,.55,11,GREY)
text(s,'2026 的变化不只在更稀疏，而在把动态 mask 的 planner 成本、训练 JVP 和 serving KV 纳入算子接口。',.82,5.5,11.5,.38,18,INK,True); foot(s,3)

# 4 taxonomy
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'多模态 mask：组合可见性，而非单一 causal')
box(s,.8,1.35,4.0,4.9,INK,INK,True); text(s,'Token stream',1.1,1.7,2.8,.3,18,WHITE,True)
for i,(l,c) in enumerate([('reasoner / text / state',MINT),('video / audio / action',GOLD),('noisy diffusion chunk',CORAL),('keyframe / reference',MINT)]): box(s,1.15,2.27+i*.74,3.25,.48,c,c,True); text(s,l,1.35,2.38+i*.74,2.75,.22,13,INK,True)
card(s,'block-causal','历史 chunk 可读；未来 chunk 禁止。',5.25,1.45,3.35,1.25,MINT)
card(s,'local + anchor','局部时空窗口加 global/keyframe bridge。',8.95,1.45,3.35,1.25,GOLD)
card(s,'read-only boundary','reasoner 不被 noisy generator 反向污染。',5.25,3.15,3.35,1.25,CORAL)
card(s,'within-chunk bidirectional','diffusion chunk 内保留去噪互动。',8.95,3.15,3.35,1.25,MINT)
text(s,'可见性条件若在 block/tile 层可判定，就传 rule / metadata；若不可判定，先 lowering 或 selector。',5.25,5.18,7.0,.48,16,INK,True); foot(s,4)

# 5 comparison
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'四种实现路径：kernel 看到的不是同一个“mask”')
cols=[('规则 / BlockMask','Causal-rCM\nFlexAttention','block schedule + predicate\n可跳过 block',MINT),('CSR / block plan','LVSA\nFlashInfer','indptr + indices\nplan 后遍历 nnz tile',GOLD),('selector / pack','VMoBA\nToken Sparse','indices + cu_seqlens\ncompact QKV',CORAL),('dense bias fallback','FrameDiT 公开代码\nSDPA/Diffusers','bool/bias broadcast\n通常仍遍历 dense tile','8C9A9C')]
for i,(h,a,b,c) in enumerate(cols):
    x=.55+i*3.18; box(s,x,1.45,2.83,4.75,WHITE,'D7DEDC',True); box(s,x,1.45,2.83,.55,c,c,True); text(s,h,x+.18,1.6,2.45,.22,12,INK,True); text(s,a,x+.2,2.33,2.38,.55,15,INK,True); text(s,b,x+.2,3.35,2.35,1.1,12,GREY)
text(s,'审查点：metadata 的大小是否为 O(nnz_blocks)？kernel grid 是否真的只遍历非零 tile？',.7,6.55,11.7,.3,15,INK,True); foot(s,5)

# 6 LVSA
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'LVSA：CPU 保留 CSR metadata，FlashInfer 消费计划')
for x,l,c in [(1.0,'geometry\nwindow + rotating anchors',MINT),(4.25,'CPU CSR\nint32 indptr / indices',GOLD),(7.5,'FlashInfer plan\nblock traversal',CORAL),(10.75,'GPU run\nskip unlisted tile',MINT)]:
    box(s,x,2.25,1.92,1.25,c,c,True); text(s,l,x+.15,2.57,1.62,.58,14,INK,True,align=PP_ALIGN.CENTER)
for x in [2.96,6.21,9.46]: text(s,'->',x,2.62,.4,.3,22,INK,True,align=PP_ALIGN.CENTER)
text(s,'源码证据：`ring_block_frame_csr` 返回 int32 CSR；`ensure_device()` 刻意不移动 fi_indptr / fi_indices，由 host builder 与 FlashInfer planning pass 消费。',.95,4.35,11.65,.55,16,INK,True)
card(s,'为何这比 dense mask 可扩展','metadata 从 O(L^2) 改为 O(nnz_blocks + n_rows)，但需计入 planner、CSR 复制与非连续 KV 访存。',1.0,5.35,11.2,.85,GOLD); foot(s,6)

# 7 Causal
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'Causal-rCM：custom mask 是 JVP operator contract')
box(s,.85,1.45,3.7,4.75,INK,INK,True); text(s,'Packed stream',1.2,1.75,2,.25,17,WHITE,True)
for i,(l,c) in enumerate([('clean block 0',MINT),('clean block 1',MINT),('noisy block 0',CORAL),('noisy block 1',CORAL)]): box(s,1.2,2.25+i*.73,2.95,.46,c,c,True); text(s,l,1.4,2.36+i*.73,2.5,.2,13,INK,True)
card(s,'BlockPattern + AttnMaskSpec','frame-token geometry、chunk、sliding window、sink、offset；不是 L x L 张量。',5.1,1.55,3.25,1.25,MINT)
card(s,'Flex BlockMask / JVP Triton','同一 teacher-forcing mask 进入 primal 与 JVP；错误的后处理 mask 不等价。',8.82,1.55,3.25,1.25,CORAL)
card(s,'KV cache + CP','同一模式覆盖 packed training、replay、inference；对齐和负载均衡仍是成本。',5.1,3.4,6.97,1.25,GOLD)
text(s,'把“mask 支持”当作 forward-only feature 会漏掉 backward/JVP、cache 与 sequence parallel 的系统合同。',5.1,5.35,6.9,.42,17,INK,True); foot(s,7)

# 8 unified
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'统一模型：先进行语义 lowering，再考虑 sparse kernel')
box(s,.8,1.55,3.4,3.55,WHITE,'D7DEDC',True); text(s,'通用 FlexAttention\n一个混合 mask',1.2,2.05,2.55,.65,19,INK,True,align=PP_ALIGN.CENTER); text(s,'正确但 kernel 对双流结构不透明\n可能做 padding-equivalent work',1.15,3.35,2.7,.55,13,GREY,align=PP_ALIGN.CENTER)
text(s,'semantic\nlowering',4.32,2.72,.8,.55,13,CORAL,True,align=PP_ALIGN.CENTER)
box(s,5.3,1.55,6.95,3.55,INK,INK,True); text(s,'Cosmos 3 two-way flat attention',5.7,1.95,5.8,.28,20,WHITE,True)
for i,(l,c) in enumerate([('reasoner causal varlen call',MINT),('generator full varlen call',GOLD)]): box(s,5.8,2.6+i*.85,5.85,.55,c,c,True); text(s,l,6.05,2.76+i*.85,5.2,.22,15,INK,True)
text(s,'本地论文材料：相对 FlexAttention baseline，Nano 训练吞吐 +22%。',.9,5.7,11.4,.34,17,INK,True); foot(s,8)

# 9 host device
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'长序列：CPU 可以规划稀疏，但不能输出 dense mask')
rows=[('静态几何','CPU 初始化 / cache','CSR / page table','异步 H2D 或 runtime plan'),('每请求 KV 选择','GPU selector 为主','selected pages / indptr','paged attention'),('每 step 小变化','reuse + delta','cached metadata','减少 planner 次数'),('每 token top-k','GPU','indices + compact QKV','避免 PCIe 往返')]
headers=['场景','生成位置','传递对象','执行要点']
for i,h in enumerate(headers): text(s,h,.8+[0,2.4,5.0,8.0][i],1.45,[2,2.2,2.8,3.7][i],.3,14,INK,True)
for r,(a,b,c,d) in enumerate(rows):
    y=2.0+r*.88; box(s,.72,y,11.9,.63,WHITE,'D8DEDC',False)
    for i,v in enumerate([a,b,c,d]): text(s,v,.88+[0,2.4,5.0,8.0][i],y+.15,[2,2.2,2.8,3.7][i],.28,12,INK if i==0 else GREY,i==0)
text(s,'禁止路径：CPU 生成 [L,L] bool/fp16 mask 再拷 GPU。它具有 O(L²) 内存、PCIe 传输与同步三重成本。',.8,6.15,11.7,.38,17,CORAL,True); foot(s,9)

# 10 checklist
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER); title(s,'实现建议：接口、指标与质量守护')
card(s,'接口','MaskSpec(kind, geometry, dynamic, storage)\nplan(spec, layout, device) -> Plan\nrun(q,k,v,Plan) -> o',.8,1.45,3.8,2.1,MINT)
card(s,'性能模型','T = T_select + T_pack + T_plan + T_attn + T_unpack\n测有效带宽、tile occupancy、nnz 曲线。',4.78,1.45,3.8,2.1,GOLD)
card(s,'正确性','dense reference、JVP/backward、chunk 边界、cache reuse；不要只验证 forward。',8.76,1.45,3.8,2.1,CORAL)
for i,(a,b) in enumerate([('视频','motion / identity / loop'),('跨模态','audio-action sync / grounding'),('服务','TTFT / TPOT / mixed batch')]):
    x=1.0+i*4.0; box(s,x,4.35,3.35,1.1,INK,INK,True); text(s,a,x+.18,4.58,.7,.23,14,WHITE,True); text(s,b,x+.86,4.58,2.2,.25,13,'C8D6D3')
foot(s,10)

# 11 conclusion
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(INK); title(s,'最终建议：让 mask 的语义在正确层级消失或变小',True)
items=[('1','能拆就拆','双流/矩形可见性 -> causal/full varlen calls',MINT),('2','需稀疏就索引化','window/anchor -> CSR / BlockMask / page table',GOLD),('3','需动态就打包','router/selector -> GPU indices + compact QKV',CORAL),('4','把控制面纳入 KPI','planner、H2D、top-k、pack/unpack 与 attention 同测',MINT)]
for i,(n,h,b,c) in enumerate(items):
    y=1.45+i*1.13; box(s,.8,y,.55,.55,c,c,True); text(s,n,.98,y+.13,.18,.18,14,INK,True,align=PP_ALIGN.CENTER); text(s,h,1.65,y+.05,2.2,.25,17,WHITE,True); text(s,b,4.0,y+.1,7.7,.3,14,'C8D6D3')
text(s,'最终文件：01_ai_infra/kernel/custom_attn/ 多模态稀疏Attention与定制Mask_Kernel调研.md / .pptx',.82,6.65,11.5,.25,11,'A9BFBA'); foot(s,11,True)

prs.save(OUT)
print(OUT)
